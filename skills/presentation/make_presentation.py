"""
Генерация презентации: Роль буддизма в ценностном пространстве Большой Евразии
Стильный дизайн с фотореалистичными иллюстрациями и логотипом Сангхи России
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ─── Color Palette ───────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x0D, 0x14)   # near-black deep blue
PANEL_BG     = RGBColor(0x16, 0x15, 0x22)   # dark panel
CARD_BG      = RGBColor(0x1E, 0x1C, 0x2E)   # card / elevated surface
GOLD         = RGBColor(0xC9, 0xA9, 0x62)   # refined gold
GOLD_BRIGHT  = RGBColor(0xE8, 0xCB, 0x7B)   # bright gold accent
WHITE        = RGBColor(0xF2, 0xF0, 0xEB)   # warm white
SUBTLE       = RGBColor(0x9A, 0x97, 0xA8)   # subtle gray text
ACCENT_LINE  = RGBColor(0x3A, 0x36, 0x52)   # muted border line

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

LOGO_PATH = "/Users/oboton/Music/Docs/Presentation_Project/logo.png"
IMG_DIR   = "/Users/oboton/Music/Docs/Presentation_Project"

# ─── Slide data ──────────────────────────────────────────────────────────────

slides_data = [
    {
        "type": "title",
        "title": "Роль буддизма\nв ценностном пространстве\nБольшой Евразии",
        "subtitle": "Буддийская Традиционная Сангха России",
    },
    {
        "type": "content",
        "title": "Исторический контекст",
        "image": "img_historical.png",
        "points": [
            "Буддизм — одна из старейших духовных традиций человечества, существующая более 2 500 лет",
            "Как мировая религия буддизм распространился по всей Азии, сформировав единое культурное пространство",
            "Общие ценности — ненасилие, сострадание, поиск истины — объединяют народы Евразии",
            "Буддийские монастыри — живые центры сохранения, изучения и передачи традиций",
        ],
    },
    {
        "type": "content_right",
        "title": "Буддизм в евразийском\nкультурном пространстве",
        "image": "img_cultural.png",
        "points": [
            "Буддийская цивилизация охватывает народы Центральной, Восточной и Юго-Восточной Азии",
            "Единое духовное пространство: общая священная литература, иконография, архитектура, этика",
            "Буддийские культурные коды — основа диалога между цивилизациями Большой Евразии",
            "Искусство, медицина, образование — плоды межкультурного буддийского обмена",
        ],
    },
    {
        "type": "content",
        "title": "Буддийская философия\nкак основа ценностного диалога",
        "image": "img_philosophy.png",
        "points": [
            "Взаимозависимость — все явления связаны между собой",
            "Сострадание — основа ненасильственного взаимодействия цивилизаций",
            "Серединный путь: отказ от крайностей как метод гармонизации отношений",
            "Стремление к просветлению ради блага всех существ",
        ],
    },
    {
        "type": "content_right",
        "title": "Этика ненасилия\nи миростроительство",
        "image": "img_nonviolence.png",
        "points": [
            "Принцип ненасилия — основополагающая ценность буддийской этики",
            "Буддийские традиции мирного диалога и современные миротворческие инициативы",
            "Притча о «раскалённом угле»: культура вражды разрушает прежде всего самого носителя",
            "Внутренняя собранность и самоконтроль — залог мира во внешних отношениях",
        ],
    },
    {
        "type": "content",
        "title": "Буддизм\nи современные вызовы",
        "image": "img_modern.png",
        "points": [
            "Высокая адаптивность буддизма к условиям цифровизации и социальных трансформаций",
            "Буддийская критика алчности: альтернатива логике бесконечного потребления",
            "Притча о «плоте»: технологии — инструмент, а не самоцель; ценностное ядро неизменно",
            "Осознанность и умеренность — буддийский ответ на экологические и духовные кризисы",
        ],
    },
    {
        "type": "content_right",
        "title": "Россия\nи евразийский буддизм",
        "image": "img_historical.png",
        "points": [
            "Буддизм — часть духовного фундамента российской идентичности наряду с православием и исламом",
            "Буддийские регионы: Бурятия, Калмыкия, Тыва — живые мосты к цивилизациям Востока",
            "Раскрытие «внутреннего Востока» России: созерцательное, сострадательное сознание",
            "Притча о горчичном зерне: за различием традиций — единый опыт человеческой уязвимости",
        ],
    },
    {
        "type": "content",
        "title": "Роль буддизма сегодня:\nключевые аспекты",
        "image": "img_today.png",
        "points": [
            "Миротворческая функция: мирное взаимодействие и взаимопроникновение культур",
            "Нравственный ориентир: этика ненасилия, сострадания и личной ответственности",
            "Интеграционный потенциал: укрепление горизонтальных связей в Большой Евразии",
            "Противовес вестернизации: сохранение культурной самобытности народов",
            "Диалог государства и духовенства: гармонизация общественных отношений",
        ],
    },
    {
        "type": "conclusion",
        "title": "Заключение",
        "image": "img_conclusion.png",
        "points": [
            "Роль буддизма в ценностном пространстве Большой Евразии сохраняется и возрастает",
            "Традиционные духовные учения помогают нащупывать пути к справедливому и устойчивому миропорядку",
            "Основа — взаимное уважение, сострадание и осознание глубинной взаимосвязанности всех существ",
            "Буддийская мудрость — ресурс для преодоления кризисов и построения гармоничного будущего",
        ],
    },
]


# ─── Helper functions ────────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, w, h, fill=None, border_color=None, border_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def rounded_rect(slide, left, top, w, h, fill=None, border_color=None, border_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, left, top, w, h, text, size=Pt(14), bold=False, italic=False,
            color=WHITE, align=PP_ALIGN.LEFT, font_name="Arial"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def multiline_textbox(slide, left, top, w, h, text, size=Pt(14), bold=False,
                      color=WHITE, align=PP_ALIGN.LEFT, line_spacing=1.15, font_name="Arial"):
    """Textbox with proper line spacing for multi-line text."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_logo(slide, left, top, size=Inches(0.55)):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, size, size)


def add_image_clipped(slide, img_path, left, top, width, height):
    """Add image maintaining aspect ratio, cropped to fill the given area."""
    if not os.path.exists(img_path):
        # fallback: just add a dark rect
        rounded_rect(slide, left, top, width, height, fill=CARD_BG, border_color=ACCENT_LINE)
        return

    pic = slide.shapes.add_picture(img_path, left, top, width, height)

    # Compute crop to fill without distortion
    img = Image.open(img_path)
    img_w, img_h = img.size
    img_aspect = img_w / img_h
    box_aspect = width / height

    if img_aspect > box_aspect:
        # Image is wider — crop sides
        crop_frac = 1.0 - (box_aspect / img_aspect)
        crop_each = crop_frac / 2.0
        pic.crop_left = crop_each
        pic.crop_right = crop_each
    else:
        # Image is taller — crop top/bottom
        crop_frac = 1.0 - (img_aspect / box_aspect)
        crop_each = crop_frac / 2.0
        pic.crop_top = crop_each
        pic.crop_bottom = crop_each

    return pic


def bullet_block(slide, points, left, top, width, point_size=Pt(14),
                 color=WHITE, spacing=Inches(0.72), dot_color=GOLD):
    """Render a list of bullet points with gold square dots."""
    y = top
    for point in points:
        # Gold dot
        rect(slide, left, y + Pt(5), Pt(7), Pt(7), fill=dot_color)
        # Text
        textbox(slide, left + Inches(0.3), y - Pt(2), width - Inches(0.3),
                Inches(0.6), point, size=point_size, color=color)
        y += spacing


def footer_bar(slide):
    """Standard footer bar with logo and site."""
    # Thin gold line
    rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Pt(1), fill=ACCENT_LINE)
    # Footer bg
    rect(slide, 0, SLIDE_H - Inches(0.52), SLIDE_W, Inches(0.52), fill=DARK_BG)
    # Logo small
    add_logo(slide, Inches(0.3), SLIDE_H - Inches(0.47), size=Inches(0.35))
    # Text
    textbox(slide, Inches(0.75), SLIDE_H - Inches(0.43), Inches(6), Inches(0.35),
            "Буддийская Традиционная Сангха России  ·  sangharussia.ru",
            size=Pt(8), color=SUBTLE, italic=True)


# ─── Slide builders ──────────────────────────────────────────────────────────

def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    # Background photo strip (right side, subtle) — Иволгинский Дацан
    add_image_clipped(slide, os.path.join(IMG_DIR, "img_russia.png"),
                      Inches(6.5), 0, Inches(6.83), SLIDE_H)

    # Dark overlay on photo (semi-transparent)
    overlay = rect(slide, Inches(6.5), 0, Inches(6.83), SLIDE_H, fill=DARK_BG)
    # Make semi-transparent via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = overlay._element.find(qn('p:spPr'))
    solidFill_el = spPr.find(qn('a:solidFill'))
    if solidFill_el is not None:
        srgb = solidFill_el.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', '55000')  # 55% opacity

    # Left panel area
    rect(slide, 0, 0, Inches(6.5), SLIDE_H, fill=DARK_BG)

    # Top thin gold line
    rect(slide, 0, Inches(0.12), Inches(6.0), Pt(2), fill=GOLD)

    # Logo
    add_logo(slide, Inches(0.6), Inches(0.7), size=Inches(1.2))

    # Org name under logo
    textbox(slide, Inches(2.0), Inches(0.85), Inches(4.0), Inches(0.6),
            "Буддийская Традиционная\nСангха России",
            size=Pt(11), color=GOLD, align=PP_ALIGN.LEFT)

    # Main title
    multiline_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.5), Inches(3.0),
                      data["title"],
                      size=Pt(28), bold=True, color=WHITE, line_spacing=1.1)

    # Gold separator
    rect(slide, Inches(0.6), Inches(5.3), Inches(4.5), Pt(2), fill=GOLD)

    # Bottom accent
    textbox(slide, Inches(0.6), Inches(6.3), Inches(5.5), Inches(0.4),
            "☸  sangharussia.ru",
            size=Pt(10), color=SUBTLE)

    # Bottom gold line
    rect(slide, 0, SLIDE_H - Inches(0.12), Inches(6.0), Pt(2), fill=GOLD)


def build_content_left_image(prs, data):
    """Image on left, text on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    img_path = os.path.join(IMG_DIR, data["image"])
    img_left = Inches(0.4)
    img_top = Inches(0.35)
    img_w = Inches(5.2)
    img_h = Inches(6.3)

    # Image card background
    rounded_rect(slide, img_left - Pt(3), img_top - Pt(3),
                 img_w + Pt(6), img_h + Pt(6),
                 fill=CARD_BG, border_color=ACCENT_LINE, border_w=Pt(1))

    # Image
    add_image_clipped(slide, img_path, img_left, img_top, img_w, img_h)

    # Right content area
    text_left = Inches(6.0)
    text_w = Inches(6.9)

    # Title
    multiline_textbox(slide, text_left, Inches(0.5), text_w, Inches(1.4),
                      data["title"],
                      size=Pt(24), bold=True, color=WHITE, line_spacing=1.05)

    # Gold line under title
    rect(slide, text_left, Inches(1.85), Inches(3.5), Pt(2), fill=GOLD)

    # Bullet points
    bullet_block(slide, data["points"],
                 left=text_left, top=Inches(2.2),
                 width=text_w, point_size=Pt(14),
                 spacing=Inches(0.85))

    # Footer
    footer_bar(slide)


def build_content_right_image(prs, data):
    """Text on left, image on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    img_path = os.path.join(IMG_DIR, data["image"])
    img_w = Inches(5.2)
    img_h = Inches(6.3)
    img_left = SLIDE_W - img_w - Inches(0.4)
    img_top = Inches(0.35)

    # Image card background
    rounded_rect(slide, img_left - Pt(3), img_top - Pt(3),
                 img_w + Pt(6), img_h + Pt(6),
                 fill=CARD_BG, border_color=ACCENT_LINE, border_w=Pt(1))

    # Image
    add_image_clipped(slide, img_path, img_left, img_top, img_w, img_h)

    # Left content area
    text_left = Inches(0.5)
    text_w = Inches(6.7)

    # Title
    multiline_textbox(slide, text_left, Inches(0.5), text_w, Inches(1.4),
                      data["title"],
                      size=Pt(24), bold=True, color=WHITE, line_spacing=1.05)

    # Gold line under title
    rect(slide, text_left, Inches(1.85), Inches(3.5), Pt(2), fill=GOLD)

    # Bullet points
    bullet_block(slide, data["points"],
                 left=text_left, top=Inches(2.2),
                 width=text_w, point_size=Pt(14),
                 spacing=Inches(0.85))

    # Footer
    footer_bar(slide)


def build_conclusion_slide(prs, data):
    """Conclusion slide — same layout as content_left_image for visual consistency."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    img_path = os.path.join(IMG_DIR, data["image"])
    img_left = Inches(0.4)
    img_top = Inches(0.35)
    img_w = Inches(5.2)
    img_h = Inches(6.3)

    # Image card background
    rounded_rect(slide, img_left - Pt(3), img_top - Pt(3),
                 img_w + Pt(6), img_h + Pt(6),
                 fill=CARD_BG, border_color=ACCENT_LINE, border_w=Pt(1))

    # Image
    add_image_clipped(slide, img_path, img_left, img_top, img_w, img_h)

    # Right content area
    text_left = Inches(6.0)
    text_w = Inches(6.9)

    # Title
    multiline_textbox(slide, text_left, Inches(0.5), text_w, Inches(1.4),
                      data["title"],
                      size=Pt(24), bold=True, color=WHITE, line_spacing=1.05)

    # Gold line under title
    rect(slide, text_left, Inches(1.45), Inches(3.5), Pt(2), fill=GOLD)

    # Bullet points
    bullet_block(slide, data["points"],
                 left=text_left, top=Inches(1.8),
                 width=text_w, point_size=Pt(14),
                 spacing=Inches(0.85))

    # Footer
    footer_bar(slide)


# ─── Build presentation ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

for slide_data in slides_data:
    stype = slide_data["type"]
    if stype == "title":
        build_title_slide(prs, slide_data)
    elif stype == "content":
        build_content_left_image(prs, slide_data)
    elif stype == "content_right":
        build_content_right_image(prs, slide_data)
    elif stype == "conclusion":
        build_conclusion_slide(prs, slide_data)

output = "/Users/oboton/Music/Docs/Буддизм_и_Большая_Евразия.pptx"
prs.save(output)
print(f"✅ Presentation saved: {output}")
