"""
Генерация презентации: Святые места Бурятии
Стильный дизайн с фотореалистичными иллюстрациями и логотипом Сангхи России
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ─── Color Palette ───────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x0D, 0x14)   # near-black deep blue
PANEL_BG     = RGBColor(0x16, 0x15, 0x22)   # dark panel
CARD_BG      = RGBColor(0x1E, 0x1C, 0x2E)   # card / elevated surface
GOLD         = RGBColor(0xC9, 0xA9, 0x62)   # refined gold
GOLD_BRIGHT  = RGBColor(0xE8, 0xCB, 0x7B)   # bright gold accent
WHITE        = RGBColor(0xF2, 0xF0, 0xEB)   # warm white
SUBTLE       = RGBColor(0x9A, 0x97, 0xA8)   # subtle gray text
ACCENT_LINE  = RGBColor(0x3A, 0x36, 0x52)   # muted border line

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

LOGO_PATH = "/Users/oboton/Music/Docs/Presentation_Project/logo.png"
IMG_DIR   = "/Users/oboton/Music/Docs/Presentation_Project"

# ─── Slide data ──────────────────────────────────────────────────────────────

slides_data = [
    {
        "type": "title",
        "title": "Святые места Бурятии —\nуникальный ресурс\nдля развития паломничества\nи религиозного туризма",
        "subtitle": "Буддийская Традиционная Сангха России",
    },
    {
        "type": "content",
        "title": "Иволгинский дацан —\nдуховная столица буддизма\nв России",
        "image": "/Users/oboton/Music/Docs/ivolginsky_photo_only_filter_png_1773867091786.png",
        "points": [
            "Главный буддийский монастырский комплекс страны — «Гандан Геже Даши Чойнхорлин», основан в 1945 году",
            "Десять дуганов-храмов, пять субурганов-ступ, оранжерея со священным деревом Бодхи",
            "Буддийский университет «Даши Чойнхорлин» — единственное учебное заведение буддийской философии в России",
            "Десятки тысяч паломников ежегодно; развитая инфраструктура — гостиница, столовая, Галерея искусств",
        ],
    },
    {
        "type": "content_right",
        "title": "Нетленное тело\nХамбо Ламы Этигэлова",
        "image": "/Users/oboton/Music/Docs/lama_bw_editorial_png_1773868082909.png",
        "points": [
            "XII Пандито Хамбо Лама Этигэлов погрузился в медитацию в 1927 году в позе лотоса",
            "В 2002 году саркофаг официально вскрыт — тело нетленно спустя 75 лет без мумификации",
            "С точки зрения естественных наук это явление необъяснимо — тело сохраняет позу без подпорок",
            "Ежегодный Праздник Хамбо Ламы Этигэлова — масштабное событие с молебном и спортивными состязаниями",
        ],
    },
    {
        "type": "content",
        "title": "Гора Бархан-Уула —\nсвященная вершина\nБаргузинской долины",
        "image": "/Users/oboton/Music/Docs/mountains_magazine_filter_png_1773868505221.png",
        "points": [
            "Одна из главных святынь Баргузинской долины, упоминаемая в древних тибетских текстах",
            "Охраняет буддийское учение с севера; место медитации великого йогина Соодой-ламы",
            "Ежегодный молебен «Барха тахилга» и паломнические восхождения на вершину горы",
            "Первое место на конкурсе «7 чудес природы Бурятии» в 2009 году (30 000+ голосов)",
        ],
    },
    {
        "type": "content_right",
        "title": "Лик богини Янжимы —\nнерукотворное чудо\nБаргузинского района",
        "image": "/Users/oboton/Music/Docs/yanzhima_magazine_filter_png_1773868771111.png",
        "points": [
            "В 2005 году на скале близ села Ярикта обнаружен нерукотворный образ танцующей богини",
            "Янжима (Сарасвати) — богиня искусств, наук, мудрости и покровительница материнства",
            "Хамбо лама Аюшеев обнаружил лик во время медитации при поисках буддийских реликвий",
            "Рядом возведён «Дворец богини Янжимы»; место взято под опеку Сангхой России",
        ],
    },
    {
        "type": "content",
        "title": "Ступа Джарун Хашор —\n«Ступа, исполняющая\nжелания»",
        "image": "/Users/oboton/Music/Docs/stupa_magazine_filter_png_1773867556689.png",
        "points": [
            "Аналог великой ступы Бодхнатх в Непале; высота 33 м, основание 44×44 м",
            "13 ступенек символизируют путь избавления от земных мук и погружения в Нирвану",
            "Первоначально построена в 1919 г., разрушена в 1937 г., восстановлена и освящена в 2001 г.",
            "В центре — учительский храм с 64 окошками и портретами, дуганы Авалокетишвары и 21-й Тары",
        ],
    },
    {
        "type": "content_right",
        "title": "Скала Шаманка\n(мыс Бурхан) —\nсвятыня Байкала",
        "image": "img_shamanka.png",
        "points": [
            "Одна из девяти святынь Азии — двухвершинная скала на острове Ольхон, озеро Байкал",
            "Место встречи шаманской и буддийской традиций; сквозная Шаманская пещера длиной 12 м",
            "Здесь совершались культовые обряды и паломничества со времён первых шаманов",
            "Государственный природно-исторический памятник; археологические находки эпохи неолита",
        ],
    },
    {
        "type": "content",
        "title": "Сакральная география\nБурятии — сеть\nпаломнических маршрутов",
        "image": "/Users/oboton/Music/Docs/buddhist_sites_collage_magazine_filter_png_1773872819730.png",
        "points": [
            "Сандаловый Будда «Зандан Жуу» — первая скульптура Будды, созданная при его жизни (высота 2,18 м)",
            "Этигэловский святой источник — комплекс целебных ключей, каждый помогает от определённых болезней",
            "Более 260 лет буддийской традиции — живая сеть дацанов, ступ и святынь мирового значения",
        ],
    },
    {
        "type": "conclusion",
        "title": "Заключение",
        "image": "/Users/oboton/Music/Docs/river_canyon_magazine_filter_png_1773871238014.png",
        "points": [
            "Святыни Бурятии — уникальный ресурс для паломничества и религиозного туризма в масштабах России",
            "Нетленное тело Хамбо Ламы Этигэлова — явление, не имеющее аналогов в мире",
            "Иволгинский дацан, сакральные горы, источники — целостное духовное пространство",
            "Сангха России готова к сотрудничеству в создании паломнических маршрутов",
        ],
    },
]


# ─── Helper functions ────────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, w, h, fill=None, border_color=None, border_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def rounded_rect(slide, left, top, w, h, fill=None, border_color=None, border_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, left, top, w, h, text, size=Pt(14), bold=False, italic=False,
            color=WHITE, align=PP_ALIGN.LEFT, font_name="Arial"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def multiline_textbox(slide, left, top, w, h, text, size=Pt(14), bold=False,
                      color=WHITE, align=PP_ALIGN.LEFT, line_spacing=1.15, font_name="Arial"):
    """Textbox with proper line spacing for multi-line text."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_logo(slide, left, top, size=Inches(0.55)):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, size, size)


def add_image_clipped(slide, img_path, left, top, width, height, round_corners=True):
    """Physically crop and round image using PIL with a solid background to fix PPTX viewers."""
    from PIL import Image, ImageDraw
    import tempfile
    import os
    if not os.path.exists(img_path):
        rounded_rect(slide, left, top, width, height, fill=CARD_BG, border_color=ACCENT_LINE)
        return

    img = Image.open(img_path).convert("RGBA")
    img_w, img_h = img.size
    box_aspect = width / height
    img_aspect = img_w / img_h

    if img_aspect > box_aspect:
        new_w = int(img_h * box_aspect)
        offset = (img_w - new_w) // 2
        img = img.crop((offset, 0, offset + new_w, img_h))
    else:
        new_h = int(img_w / box_aspect)
        offset = (img_h - new_h) // 2
        img = img.crop((0, offset, img_w, offset + new_h))

    target_w = int(width.inches * 144)
    target_h = int(height.inches * 144)
    
    try:
        resample_filter = Image.Resampling.LANCZOS
    except AttributeError:
        resample_filter = Image.ANTIALIAS
        
    img = img.resize((target_w, target_h), resample_filter)

    if round_corners:
        mask = Image.new("L", (target_w, target_h), 0)
        draw = ImageDraw.Draw(mask)
        rad = int(min(target_w, target_h) * 0.04)
        draw.rounded_rectangle((0, 0, target_w, target_h), radius=rad, fill=255)
        
        # Create a solid background matching CARD_BG
        # CARD_BG is an RGBColor object, we can read its tuple
        bg_color = CARD_BG
        if hasattr(bg_color, 'rgb'):
            bg_color = CARD_BG.rgb
        if type(bg_color) == tuple:
            r, g, b = bg_color
        else:
            # hex tuple
            hx = str(bg_color)
            r, g, b = int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)
            
        bg = Image.new("RGB", (target_w, target_h), (r, g, b))
        bg.paste(img, (0, 0), mask=mask)
        img = bg
    else:
        img = img.convert("RGB")

    temp_path = os.path.join(tempfile.gettempdir(), "cropped_img_temp_" + str(hash(img_path)) + ".jpg")
    img.save(temp_path, "JPEG", quality=95)

    pic = slide.shapes.add_picture(temp_path, left, top, width, height)
    return pic


def bullet_block(slide, points, left, top, width, point_size=Pt(14),
                 color=WHITE, spacing=Inches(0.72), dot_color=GOLD):
    """Render a list of bullet points with gold square dots."""
    y = top
    for point in points:
        # Gold dot
        rect(slide, left, y + Pt(5), Pt(7), Pt(7), fill=dot_color)
        # Text
        textbox(slide, left + Inches(0.3), y - Pt(2), width - Inches(0.3),
                Inches(0.6), point, size=point_size, color=color)
        y += spacing


def footer_bar(slide):
    """Standard footer bar with logo and site."""
    # Thin gold line
    rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Pt(1), fill=ACCENT_LINE)
    # Footer bg
    rect(slide, 0, SLIDE_H - Inches(0.52), SLIDE_W, Inches(0.52), fill=DARK_BG)
    # Logo small
    add_logo(slide, Inches(0.3), SLIDE_H - Inches(0.47), size=Inches(0.35))
    # Text
    textbox(slide, Inches(0.75), SLIDE_H - Inches(0.43), Inches(6), Inches(0.35),
            "Буддийская Традиционная Сангха России  ·  sangharussia.ru",
            size=Pt(8), color=SUBTLE, italic=True)


# ─── Slide builders ──────────────────────────────────────────────────────────

def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    # Background photo strip (right side, subtle) — Иволгинский Дацан
    add_image_clipped(slide, "/Users/oboton/Music/Docs/buddha_rock_magazine_filter_png_1773869713360.png",
                      Inches(6.5), 0, Inches(6.83), SLIDE_H, round_corners=False)

    # Dark overlay on photo (semi-transparent)
    overlay = rect(slide, Inches(6.5), 0, Inches(6.83), SLIDE_H, fill=DARK_BG)
    # Make semi-transparent via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = overlay._element.find(qn('p:spPr'))
    solidFill_el = spPr.find(qn('a:solidFill'))
    if solidFill_el is not None:
        srgb = solidFill_el.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', '55000')  # 55% opacity

    # Left panel area
    rect(slide, 0, 0, Inches(6.5), SLIDE_H, fill=DARK_BG)

    # Top thin gold line
    rect(slide, 0, Inches(0.12), Inches(6.0), Pt(2), fill=GOLD)

    # Logo
    add_logo(slide, Inches(0.6), Inches(0.7), size=Inches(1.2))

    # Org name under logo
    textbox(slide, Inches(2.0), Inches(0.85), Inches(4.0), Inches(0.6),
            "Буддийская Традиционная\nСангха России",
            size=Pt(11), color=GOLD, align=PP_ALIGN.LEFT)

    # Main title
    multiline_textbox(slide, Inches(0.6), Inches(2.3), Inches(5.5), Inches(3.5),
                      data["title"],
                      size=Pt(26), bold=True, color=WHITE, line_spacing=1.1)

    # Gold separator
    rect(slide, Inches(0.6), Inches(5.5), Inches(4.5), Pt(2), fill=GOLD)

    # Bottom accent
    textbox(slide, Inches(0.6), Inches(6.3), Inches(5.5), Inches(0.4),
            "☸  sangharussia.ru",
            size=Pt(10), color=SUBTLE)

    # Bottom gold line
    rect(slide, 0, SLIDE_H - Inches(0.12), Inches(6.0), Pt(2), fill=GOLD)


def build_content_left_image(prs, data):
    """Image on left, text on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    img_path = os.path.join(IMG_DIR, data["image"])
    img_left = Inches(0.4)
    img_top = Inches(0.35)
    img_w = Inches(5.2)
    img_h = Inches(6.3)

    # Image card background
    rounded_rect(slide, img_left - Pt(3), img_top - Pt(3),
                 img_w + Pt(6), img_h + Pt(6),
                 fill=CARD_BG, border_color=ACCENT_LINE, border_w=Pt(1))

    # Image
    add_image_clipped(slide, img_path, img_left, img_top, img_w, img_h)

    # Right content area
    text_left = Inches(6.0)
    text_w = Inches(6.9)

    # Title
    multiline_textbox(slide, text_left, Inches(0.5), text_w, Inches(1.4),
                      data["title"],
                      size=Pt(24), bold=True, color=WHITE, line_spacing=1.05)

    # Gold line under title
    rect(slide, text_left, Inches(1.85), Inches(3.5), Pt(2), fill=GOLD)

    # Bullet points
    bullet_block(slide, data["points"],
                 left=text_left, top=Inches(2.2),
                 width=text_w, point_size=Pt(14),
                 spacing=Inches(0.85))

    # Footer
    footer_bar(slide)


def build_content_right_image(prs, data):
    """Text on left, image on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    img_path = os.path.join(IMG_DIR, data["image"])
    img_w = Inches(5.2)
    img_h = Inches(6.3)
    img_left = SLIDE_W - img_w - Inches(0.4)
    img_top = Inches(0.35)

    # Image card background
    rounded_rect(slide, img_left - Pt(3), img_top - Pt(3),
                 img_w + Pt(6), img_h + Pt(6),
                 fill=CARD_BG, border_color=ACCENT_LINE, border_w=Pt(1))

    # Image
    add_image_clipped(slide, img_path, img_left, img_top, img_w, img_h)

    # Left content area
    text_left = Inches(0.5)
    text_w = Inches(6.7)

    # Title
    multiline_textbox(slide, text_left, Inches(0.5), text_w, Inches(1.4),
                      data["title"],
                      size=Pt(24), bold=True, color=WHITE, line_spacing=1.05)

    # Gold line under title
    rect(slide, text_left, Inches(1.85), Inches(3.5), Pt(2), fill=GOLD)

    # Bullet points
    bullet_block(slide, data["points"],
                 left=text_left, top=Inches(2.2),
                 width=text_w, point_size=Pt(14),
                 spacing=Inches(0.85))

    # Footer
    footer_bar(slide)


def build_conclusion_slide(prs, data):
    """Conclusion slide — same layout as content_left_image for visual consistency."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    img_path = os.path.join(IMG_DIR, data["image"])
    img_left = Inches(0.4)
    img_top = Inches(0.35)
    img_w = Inches(5.2)
    img_h = Inches(6.3)

    # Image card background
    rounded_rect(slide, img_left - Pt(3), img_top - Pt(3),
                 img_w + Pt(6), img_h + Pt(6),
                 fill=CARD_BG, border_color=ACCENT_LINE, border_w=Pt(1))

    # Image
    add_image_clipped(slide, img_path, img_left, img_top, img_w, img_h)

    # Right content area
    text_left = Inches(6.0)
    text_w = Inches(6.9)

    # Title
    multiline_textbox(slide, text_left, Inches(0.5), text_w, Inches(1.4),
                      data["title"],
                      size=Pt(24), bold=True, color=WHITE, line_spacing=1.05)

    # Gold line under title
    rect(slide, text_left, Inches(1.45), Inches(3.5), Pt(2), fill=GOLD)

    # Bullet points
    bullet_block(slide, data["points"],
                 left=text_left, top=Inches(1.8),
                 width=text_w, point_size=Pt(14),
                 spacing=Inches(0.85))

    # Footer
    footer_bar(slide)


# ─── Build presentation ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

for slide_data in slides_data:
    stype = slide_data["type"]
    if stype == "title":
        build_title_slide(prs, slide_data)
    elif stype == "content":
        build_content_left_image(prs, slide_data)
    elif stype == "content_right":
        build_content_right_image(prs, slide_data)
    elif stype == "conclusion":
        build_conclusion_slide(prs, slide_data)

output = "/Users/oboton/Music/Docs/Святые_места_Бурятии.pptx"
prs.save(output)
print(f"✅ Presentation saved: {output}")
